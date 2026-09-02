#!/usr/bin/env python3
"""Assemble the Round 2 deck on the official Accenture AIC template.

Keeps the template's cover, team-details, video and thank-you slides and places
ten designed content slides between them as full-bleed images — so the deck
carries Accenture's chrome and none of our typefaces can substitute on a machine
that lacks them.

Slides are REUSED rather than removed wherever possible: python-pptx drops a
slide from the id list but leaves its part in the package, and a later add_slide
then reuses that part name and writes it twice. The one slide that must go (the
template's own "remove before submission" instructions) is deleted afterwards by
rewriting the package, orphans and all.
"""
import os, re, shutil, zipfile
from pptx import Presentation

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TPL = os.path.expanduser("~/Downloads/6a7c763e089e4_aic_2026/AIC_Talent-Brand_PPT-Template (1).pptx")
OUT = os.path.join(ROOT, "submission", "KAIROS_Detailed-Business-Proposal.pptx")
TMP = OUT + ".stage"
DECK = os.path.join(ROOT, "submission", "deck")

pngs = sorted(os.path.join(DECK, f) for f in os.listdir(DECK) if f.endswith(".png"))
assert len(pngs) == 10, "expected 10 content slides, found %d" % len(pngs)

prs = Presentation(TPL)
W, H = prs.slide_width, prs.slide_height
LAYOUT = prs.slides[3].slide_layout                      # 1_Standard slide_no bullets

def fill(slide, png):
    for shp in list(slide.shapes):
        shp._element.getparent().remove(shp._element)
    slide.shapes.add_picture(png, 0, 0, width=W, height=H)

# reuse the two placeholder content slides, then append the remaining eight
fill(prs.slides[3], pngs[0])
fill(prs.slides[4], pngs[1])
for png in pngs[2:]:
    s = prs.slides.add_slide(LAYOUT)
    fill(s, png)

# order: cover · instructions(removed later) · team · content x10 · video · thanks
lst = prs.slides._sldIdLst
ids = list(lst)
cover, instr, team, c1, c2, video, thanks = ids[0], ids[1], ids[2], ids[3], ids[4], ids[5], ids[6]
rest = ids[7:]
for e in ids:
    lst.remove(e)
for e in [cover, instr, team, c1, c2] + rest + [video, thanks]:
    lst.append(e)

# ── team details: a single, individual entry ────────────────────────────────
# The template ships three member blocks. An individual entrant keeps one and the
# other two are removed entirely — leaving empty frames behind reads as an
# unfinished slide, not as a solo entry.
team = prs.slides[2]
by_name = {}
for shp in team.shapes:
    by_name.setdefault(shp.name, []).append(shp)

def kill(shape):
    shape._element.getparent().remove(shape._element)

# member 2 (lower-left) and member 3 (right), plus the divider that separated the columns
for shp in list(team.shapes):
    L, T = shp.left / 914400.0, shp.top / 914400.0
    drop = (
        L > 6.6                                   # everything in the right-hand column
        or (T > 4.4 and shp.shape_type != 1 and "Rectangle 2" not in shp.name)  # lower-left block
        or (shp.name == "Rectangle 29")           # the second "Photo" frame
        or (shp.name == "Straight Connector 5")   # the column divider
    )
    if shp.name in ("Title 17", "Slide Number Placeholder 9", "Rectangle 2"):
        drop = False
    if drop:
        kill(shp)

FIELDS = {
    "Name (Team Leader)": "Aswanth",
    "College:": "College: Indian Institute of Technology Madras",
    "Stream:": "Stream: Civil Engineering (B.Tech)",
    "Year of graduation:": "Year of graduation: 2028",
}
for shp in team.shapes:
    if not shp.has_text_frame:
        continue
    for para in shp.text_frame.paragraphs:
        for run in para.runs:
            key = run.text.strip()
            if key in FIELDS:
                run.text = FIELDS[key]

for shp in team.shapes:
    if shp.has_table:
        c = shp.table.cell(0, 0)
        for para in c.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip() == "TEAM NAME:":
                    run.text = "TEAM NAME:  Aswanth  (individual participant)"

# the photo, if one has been saved next to the deck
PHOTO = os.path.join(ROOT, "submission", "photo.jpg")
if not os.path.exists(PHOTO):
    for alt in ("photo.jpeg", "photo.png"):
        cand = os.path.join(ROOT, "submission", alt)
        if os.path.exists(cand):
            PHOTO = cand
            break
if os.path.exists(PHOTO):
    frame = next((sh for sh in team.shapes if sh.shape_type == 13), None)
    if frame is not None:
        L, T, W, H = frame.left, frame.top, frame.width, frame.height
        kill(frame)
        from PIL import Image
        iw, ih = Image.open(PHOTO).size
        scale = min(W / iw, H / ih)              # fit inside, keep the aspect ratio
        w, h = int(iw * scale), int(ih * scale)
        team.shapes.add_picture(PHOTO, L + (W - w) // 2, T + (H - h) // 2, width=w, height=h)
        print("photo inserted from", os.path.basename(PHOTO))
else:
    print("NO PHOTO FOUND — save it as submission/photo.jpg and re-run to insert it")

cp = prs.core_properties
cp.title = "KAIRÓS — Detailed Business Proposal"
cp.subject = "Accenture Innovation Challenge 2026 · Round 2 · BusinessIntelligence.ai"
cp.author = ""; cp.last_modified_by = ""; cp.comments = ""; cp.keywords = ""
prs.save(TMP)

# ── drop the instructions slide by rewriting the package ────────────────────
zin = zipfile.ZipFile(TMP)
pres = zin.read("ppt/presentation.xml").decode("utf8")
rels = zin.read("ppt/_rels/presentation.xml.rels").decode("utf8")

# the instructions slide is the second sldId in the list
sldids = re.findall(r'<p:sldId[^>]*/>', pres)
target = sldids[1]
rid = re.search(r'r:id="([^"]+)"', target).group(1)
tgt = re.search(r'Id="%s"[^>]*Target="([^"]+)"' % rid, rels).group(1)
drop_part = "ppt/" + tgt.replace("../", "")
drop_rels = drop_part.replace("ppt/slides/", "ppt/slides/_rels/") + ".rels"

pres = pres.replace(target, "")
rels = re.sub(r'<Relationship Id="%s"[^>]*/>' % rid, "", rels)
ct = zin.read("[Content_Types].xml").decode("utf8")
ct = re.sub(r'<Override PartName="/%s"[^>]*/>' % re.escape(drop_part), "", ct)

# every part the removed slide owned and nothing else references
keep = [n for n in zin.namelist() if n not in (drop_part, drop_rels)]
seen = set()
with zipfile.ZipFile(OUT, "w", zipfile.ZIP_DEFLATED) as zout:
    for n in keep:
        if n in seen:            # never write a part twice
            continue
        seen.add(n)
        if n == "ppt/presentation.xml":
            zout.writestr(n, pres)
        elif n == "ppt/_rels/presentation.xml.rels":
            zout.writestr(n, rels)
        elif n == "[Content_Types].xml":
            zout.writestr(n, ct)
        else:
            zout.writestr(n, zin.read(n))
zin.close(); os.remove(TMP)

# ── font remap ──────────────────────────────────────────────────────────────
# The template is set in Graphik, Accenture's brand face. It is not installed on
# a typical machine, so PowerPoint substitutes something unpredictable and warns
# the reader that fonts are missing. The template's own instruction slide says to
# use standard Arial, so we map it explicitly: a known substitution beats an
# unknown one, and the warning goes away. Our content slides are images and are
# unaffected either way.
# Graphik is the brand face; the rest are display faces the template carries on
# layouts we do not use. All of them would warn. The CJK/Indic entries below are
# left alone - they are the standard Office script fallbacks and warn about
# nothing.
GRAPHIK = ("Graphik Semibold", "Graphik Regular", "Graphik Black",
           "Graphik Extralight", "Graphik Medium", "Graphik-Semibold", "Graphik",
           "GT Sectra Fine Rg", "Gotham Medium", "Roboto Light", "System Font",
           "Aptos Display", "Aptos", "Helvetica Neue Medium", "Helvetica Neue Light",
           "Helvetica Neue", "Calibri Light")
import zipfile as _zf
_tmp = OUT + ".fонts"
_zin = _zf.ZipFile(OUT)
with _zf.ZipFile(_tmp, "w", _zf.ZIP_DEFLATED) as _zout:
    for _n in _zin.namelist():
        _d = _zin.read(_n)
        if _n.endswith(".xml") and (_n.startswith("ppt/slides/")
                                    or _n.startswith("ppt/slideLayouts/")
                                    or _n.startswith("ppt/slideMasters/")
                                    or _n.startswith("ppt/theme/")):
            _t = _d.decode("utf8")
            for _g in GRAPHIK:
                _t = _t.replace('typeface="%s"' % _g, 'typeface="Arial"')
            _d = _t.encode("utf8")
        _zout.writestr(_n, _d)
_zin.close()
os.replace(_tmp, OUT)

check = Presentation(OUT)
z = zipfile.ZipFile(OUT)
import collections
dups = [n for n, c in collections.Counter(z.namelist()).items() if c > 1]
print("saved: %s  (%.0f KB)" % (OUT, os.path.getsize(OUT) / 1024))
print("slides: %d   duplicate parts: %s" % (len(check.slides), dups or "none"))
for i, s in enumerate(check.slides, 1):
    t = [sh.text_frame.text.replace("\n", " / ")[:40]
         for sh in s.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
    print("  %2d  %-8s %s" % (i, "image" if not t else "template", " | ".join(t)[:60]))
